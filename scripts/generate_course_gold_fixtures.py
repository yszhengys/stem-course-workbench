#!/usr/bin/env python3
"""Generate the repository-owned CC0 PDF/PPTX evidence gold set.

The fixture text and diagrams are original to this repository. Generation is
offline and deterministic for the dependency versions pinned by ``uv.lock``.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import tempfile
from datetime import UTC, datetime
from pathlib import Path
from typing import Any
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPOSITORY_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = REPOSITORY_ROOT / "tests" / "course" / "fixtures" / "gold"
FIXED_TIME = datetime(2026, 8, 29, tzinfo=UTC)
ZIP_TIME = (2026, 8, 29, 0, 0, 0)
PDF_NAME = "stem-evidence-gold.pdf"
PPTX_NAME = "stem-evidence-gold.pptx"


def _font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    return ImageFont.load_default(size=size)


def _draw_arrow(
    drawing: ImageDraw.ImageDraw,
    start: tuple[int, int],
    end: tuple[int, int],
    *,
    color: str,
    width: int = 8,
) -> None:
    drawing.line((*start, *end), fill=color, width=width)
    x, y = end
    drawing.polygon([(x, y), (x - 28, y - 18), (x - 28, y + 18)], fill=color)


def _pdf_pages() -> list[Image.Image]:
    title = _font(56)
    body = _font(42)
    label = _font(34)
    pages: list[Image.Image] = []

    first = Image.new("RGB", (1600, 1200), "white")
    draw = ImageDraw.Draw(first)
    draw.text((90, 70), "Constant acceleration evidence", font=title, fill="#101828")
    draw.text((90, 170), "Formula: v(t) = v0 + a*t", font=body, fill="#101828")
    draw.text(
        (90, 245),
        "Given v0 = 2 m/s, a = 5 m/s^2, and t = 2 s.",
        font=label,
        fill="#344054",
    )
    draw.line((180, 980, 180, 390), fill="#101828", width=7)
    draw.line((180, 980, 1390, 980), fill="#101828", width=7)
    draw.line((180, 900, 1220, 480), fill="#1570ef", width=12)
    _draw_arrow(draw, (180, 900), (1220, 480), color="#1570ef", width=12)
    draw.text((1270, 955), "t", font=label, fill="#101828")
    draw.text((120, 350), "v", font=label, fill="#101828")
    draw.text((820, 530), "velocity-time graph", font=label, fill="#1570ef")
    draw.rounded_rectangle((900, 710, 1460, 870), radius=24, outline="#027a48", width=7)
    draw.text((945, 755), "ANSWER: 12 m/s", font=body, fill="#027a48")
    pages.append(first)

    second = Image.new("RGB", (1600, 1200), "white")
    draw = ImageDraw.Draw(second)
    draw.text((90, 70), "Projectile-motion question", font=title, fill="#101828")
    draw.text(
        (90, 170),
        "A ball is launched horizontally at 8 m/s from height 20 m.",
        font=label,
        fill="#344054",
    )
    draw.text(
        (90, 230),
        "Question: identify the horizontal and vertical velocity components.",
        font=label,
        fill="#344054",
    )
    draw.line((180, 980, 180, 390), fill="#101828", width=7)
    draw.line((180, 980, 1390, 980), fill="#101828", width=7)
    trajectory = [(220 + index * 50, 430 + int(0.9 * index * index)) for index in range(20)]
    draw.line(trajectory, fill="#b42318", width=12)
    _draw_arrow(draw, (220, 430), (570, 430), color="#1570ef", width=10)
    _draw_arrow(draw, (570, 430), (570, 760), color="#f79009", width=10)
    draw.text((320, 365), "vx = 8 m/s", font=label, fill="#1570ef")
    draw.text((600, 570), "vy changes", font=label, fill="#b54708")
    draw.text((900, 880), "x-y component diagram", font=label, fill="#b42318")
    pages.append(second)

    return pages


def _write_pdf(path: Path) -> None:
    pages = _pdf_pages()
    pages[0].save(
        path,
        "PDF",
        save_all=True,
        append_images=pages[1:],
        resolution=144,
        title="STEM Course Workbench evidence gold set",
        author="STEM Course Workbench contributors",
        subject="CC0 test fixture",
        creationDate=FIXED_TIME.timetuple(),
        modDate=FIXED_TIME.timetuple(),
    )


def _set_text(shape: Any, text: str, *, size: int, color: str) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Arial"
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = RGBColor.from_string(color)


def _add_title(slide: Any, text: str) -> None:
    title = slide.shapes.add_textbox(Inches(0.65), Inches(0.35), Inches(12), Inches(0.7))
    _set_text(title, text, size=28, color="101828")


def _write_pptx_uncanonicalized(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = "STEM Course Workbench evidence gold set"
    presentation.core_properties.subject = "CC0 test fixture"
    presentation.core_properties.author = "STEM Course Workbench contributors"
    presentation.core_properties.last_modified_by = "STEM Course Workbench contributors"
    presentation.core_properties.created = FIXED_TIME.replace(tzinfo=None)
    presentation.core_properties.modified = FIXED_TIME.replace(tzinfo=None)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_title(slide, "Vector diagram: low-text visual evidence")
    origin_x, origin_y = Inches(2.0), Inches(5.8)
    horizontal = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, origin_x, origin_y, Inches(8.5), origin_y
    )
    horizontal.line.color.rgb = RGBColor(21, 112, 239)
    horizontal.line.width = Pt(7)
    horizontal.line.end_arrowhead = True
    vertical = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, origin_x, origin_y, origin_x, Inches(1.7)
    )
    vertical.line.color.rgb = RGBColor(247, 144, 9)
    vertical.line.width = Pt(7)
    vertical.line.end_arrowhead = True
    vector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, origin_x, origin_y, Inches(8.5), Inches(1.7)
    )
    vector.line.color.rgb = RGBColor(180, 35, 24)
    vector.line.width = Pt(10)
    vector.line.end_arrowhead = True
    vector_label = slide.shapes.add_textbox(Inches(5.0), Inches(3.0), Inches(2.2), Inches(0.6))
    _set_text(vector_label, "vector v", size=24, color="B42318")

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_title(slide, "Coordinate components and constant acceleration")
    formula = slide.shapes.add_textbox(Inches(1.0), Inches(1.55), Inches(11.2), Inches(1.0))
    _set_text(formula, "v(t) = v0 + a*t", size=38, color="1570EF")
    formula.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    axes = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.4), Inches(3.0), Inches(10.5), Inches(3.4))
    axes.fill.background()
    axes.line.color.rgb = RGBColor(208, 213, 221)
    component = slide.shapes.add_textbox(Inches(2.0), Inches(3.8), Inches(9.2), Inches(1.4))
    _set_text(
        component,
        "horizontal: vx = 8 m/s\nvertical: vy changes with gravity",
        size=25,
        color="344054",
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_title(slide, "Worked answer")
    equation = slide.shapes.add_textbox(Inches(1.1), Inches(1.7), Inches(11.1), Inches(1.0))
    _set_text(equation, "v = 2 m/s + (5 m/s^2)(2 s)", size=34, color="344054")
    answer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.1), Inches(3.4), Inches(7.2), Inches(1.7))
    answer.fill.solid()
    answer.fill.fore_color.rgb = RGBColor(236, 253, 243)
    answer.line.color.rgb = RGBColor(2, 122, 72)
    answer.line.width = Pt(4)
    _set_text(answer, "ANSWER: 12 m/s", size=40, color="027A48")
    answer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    presentation.save(str(path))


def _canonicalize_pptx(source: Path, destination: Path) -> None:
    with ZipFile(source) as input_archive, ZipFile(
        destination, "w", compression=ZIP_DEFLATED, compresslevel=9
    ) as output_archive:
        for name in sorted(input_archive.namelist()):
            source_info = input_archive.getinfo(name)
            info = ZipInfo(filename=name, date_time=ZIP_TIME)
            info.compress_type = ZIP_DEFLATED
            info.external_attr = source_info.external_attr
            info.create_system = 3
            output_archive.writestr(info, input_archive.read(name))


def _write_pptx(path: Path) -> None:
    with tempfile.TemporaryDirectory(prefix="stem-course-gold-") as temporary:
        raw = Path(temporary) / "raw.pptx"
        _write_pptx_uncanonicalized(raw)
        _canonicalize_pptx(raw, path)


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _manifest(output: Path) -> dict[str, object]:
    return {
        "fixture_version": 1,
        "license": "CC0-1.0",
        "files": [
            {
                "path": PDF_NAME,
                "sha256": _sha256(output / PDF_NAME),
                "kind": "pdf",
                "page_count": 2,
                "expected": [
                    {
                        "index": 1,
                        "text": "v(t) = v0 + a*t",
                        "category": "formula",
                        "bbox_required": True,
                    },
                    {
                        "index": 1,
                        "text": "ANSWER: 12 m/s",
                        "category": "answer",
                        "bbox_required": True,
                    },
                    {
                        "index": 2,
                        "text": "Projectile-motion question",
                        "category": "question",
                        "bbox_required": True,
                    },
                    {
                        "index": 2,
                        "text": "x-y component diagram",
                        "category": "diagram",
                        "bbox_required": True,
                    },
                ],
            },
            {
                "path": PPTX_NAME,
                "sha256": _sha256(output / PPTX_NAME),
                "kind": "pptx",
                "page_count": 3,
                "expected": [
                    {
                        "index": 1,
                        "text": "Vector diagram",
                        "category": "diagram",
                        "bbox_required": True,
                    },
                    {
                        "index": 2,
                        "text": "v(t) = v0 + a*t",
                        "category": "formula",
                        "bbox_required": True,
                    },
                    {
                        "index": 3,
                        "text": "ANSWER: 12 m/s",
                        "category": "answer",
                        "bbox_required": True,
                    },
                ],
            },
        ],
    }


def generate(output: Path) -> None:
    output.mkdir(parents=True, exist_ok=True)
    _write_pdf(output / PDF_NAME)
    _write_pptx(output / PPTX_NAME)
    manifest = _manifest(output)
    (output / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path)
    arguments = parser.parse_args()
    output = arguments.output.resolve() if arguments.output else DEFAULT_OUTPUT.resolve()
    if arguments.output is None and not output.is_relative_to(REPOSITORY_ROOT):
        raise SystemExit("default output escaped the repository")
    generate(output)


if __name__ == "__main__":
    main()
