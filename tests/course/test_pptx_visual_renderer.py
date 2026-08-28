from __future__ import annotations

import hashlib
import os
import struct
import subprocess
from pathlib import Path
from typing import Any
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from pypdf import PdfWriter

from open_notebook.course.pptx_visual_renderer import (
    PptxVisualRejected,
    PptxVisualRenderer,
    PptxVisualUnavailable,
)

GOLD_ROOT = Path(__file__).parent / "fixtures" / "gold"


def _pptx(path: Path, slide_count: int = 2) -> None:
    from pptx import Presentation

    presentation = Presentation()
    for index in range(slide_count):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(0, 0, 2_000_000, 1_000_000)
        shape.text = f"Synthetic slide {index + 1}"
    presentation.save(str(path))


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


class _FakeProcess:
    def __init__(
        self,
        arguments: list[str],
        *,
        page_count: int,
        timeout: bool = False,
        mutate_source: Path | None = None,
    ) -> None:
        self.arguments = arguments
        self.page_count = page_count
        self.timeout = timeout
        self.mutate_source = mutate_source
        self.pid = 991_337
        self.returncode: int | None = None

    def communicate(self, timeout: float) -> tuple[bytes, bytes]:
        if self.timeout:
            raise subprocess.TimeoutExpired(self.arguments, timeout)
        if self.mutate_source is not None:
            self.mutate_source.write_bytes(self.mutate_source.read_bytes() + b"changed")
        output_dir = Path(self.arguments[self.arguments.index("--outdir") + 1])
        input_path = Path(self.arguments[-1])
        writer = PdfWriter()
        for _ in range(self.page_count):
            writer.add_blank_page(width=720, height=405)
        with (output_dir / f"{input_path.stem}.pdf").open("wb") as handle:
            writer.write(handle)
        self.returncode = 0
        return b"converted", b""

    def wait(self, timeout: float) -> int:
        del timeout
        self.returncode = -15
        return self.returncode

    def kill(self) -> None:
        self.returncode = -9


class _PopenFactory:
    def __init__(
        self,
        *,
        page_count: int,
        timeout: bool = False,
        mutate_source: Path | None = None,
    ) -> None:
        self.page_count = page_count
        self.timeout = timeout
        self.mutate_source = mutate_source
        self.calls: list[tuple[list[str], dict[str, Any]]] = []

    def __call__(self, arguments: list[str], **kwargs: Any) -> _FakeProcess:
        self.calls.append((arguments, kwargs))
        return _FakeProcess(
            arguments,
            page_count=self.page_count,
            timeout=self.timeout,
            mutate_source=self.mutate_source,
        )


def _renderer(factory: _PopenFactory, **kwargs: Any) -> PptxVisualRenderer:
    return PptxVisualRenderer(
        converter=Path("/opt/safe/soffice"),
        popen_factory=factory,
        **kwargs,
    )


def test_render_uses_argument_array_and_writes_bounded_opaque_pngs(
    tmp_path: Path,
) -> None:
    source = tmp_path / "lesson with spaces.pptx"
    output = tmp_path / "visuals"
    _pptx(source, slide_count=2)
    factory = _PopenFactory(page_count=2)

    rendered = _renderer(factory).render(source, _sha256(source), output)

    assert list(rendered) == [1, 2]
    assert [path.name for path in rendered.values()] == [
        "slide-0001.png",
        "slide-0002.png",
    ]
    for path in rendered.values():
        content = path.read_bytes()
        assert content.startswith(b"\x89PNG\r\n\x1a\n")
        width, height = struct.unpack(">II", content[16:24])
        assert width == 1280
        assert 1 <= height <= 10_000

    assert len(factory.calls) == 1
    arguments, options = factory.calls[0]
    assert isinstance(arguments, list)
    assert arguments[:6] == [
        "/opt/safe/soffice",
        "--headless",
        "--nologo",
        "--nodefault",
        "--nofirststartwizard",
        "--norestore",
    ]
    assert "--convert-to" in arguments
    assert arguments[arguments.index("--convert-to") + 1] == "pdf"
    assert options["start_new_session"] is True
    assert options["stdin"] is subprocess.DEVNULL
    assert options["stdout"] is subprocess.PIPE
    assert options["stderr"] is subprocess.PIPE
    assert set(options["env"]) == {"HOME", "LANG", "LC_ALL", "PATH", "TMPDIR"}
    assert "shell" not in options


def test_render_is_deterministic_across_repeated_calls(tmp_path: Path) -> None:
    source = tmp_path / "lesson.pptx"
    _pptx(source, slide_count=2)
    renderer = _renderer(_PopenFactory(page_count=2))

    first = renderer.render(source, _sha256(source), tmp_path / "first")
    second = renderer.render(source, _sha256(source), tmp_path / "second")

    assert [path.read_bytes() for path in first.values()] == [
        path.read_bytes() for path in second.values()
    ]


@pytest.mark.parametrize(
    "member",
    [
        "../outside.xml",
        "ppt/vbaProject.bin",
        "ppt/activeX/activeX1.bin",
        "ppt/embeddings/oleObject1.bin",
    ],
)
def test_render_rejects_unsafe_zip_members(tmp_path: Path, member: str) -> None:
    source = tmp_path / "unsafe.pptx"
    _pptx(source)
    with ZipFile(source, "a", ZIP_DEFLATED) as archive:
        archive.writestr(member, b"unsafe")

    with pytest.raises(PptxVisualRejected, match="unsafe"):
        _renderer(_PopenFactory(page_count=2)).render(
            source, _sha256(source), tmp_path / "visuals"
        )


def test_render_rejects_external_relationships(tmp_path: Path) -> None:
    source = tmp_path / "external.pptx"
    _pptx(source)
    relationship = b"""<?xml version="1.0" encoding="UTF-8"?>
    <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
      <Relationship Id="rId1" Type="urn:test" Target="https://example.invalid/a" TargetMode="External"/>
    </Relationships>"""
    with ZipFile(source, "a", ZIP_DEFLATED) as archive:
        archive.writestr("ppt/slides/_rels/unsafe.xml.rels", relationship)

    with pytest.raises(PptxVisualRejected, match="external"):
        _renderer(_PopenFactory(page_count=2)).render(
            source, _sha256(source), tmp_path / "visuals"
        )


def test_render_rejects_source_hash_change_before_or_during_render(
    tmp_path: Path,
) -> None:
    source = tmp_path / "lesson.pptx"
    _pptx(source)
    original_hash = _sha256(source)

    with pytest.raises(PptxVisualRejected, match="hash"):
        _renderer(_PopenFactory(page_count=2)).render(
            source, "0" * 64, tmp_path / "before"
        )

    with pytest.raises(PptxVisualRejected, match="changed"):
        _renderer(_PopenFactory(page_count=2, mutate_source=source)).render(
            source, original_hash, tmp_path / "during"
        )


def test_render_rejects_symlink_output_and_wrong_page_count(tmp_path: Path) -> None:
    source = tmp_path / "lesson.pptx"
    _pptx(source, slide_count=2)
    real_output = tmp_path / "real"
    real_output.mkdir()
    symlink_output = tmp_path / "linked"
    symlink_output.symlink_to(real_output, target_is_directory=True)

    with pytest.raises(PptxVisualRejected, match="symlink"):
        _renderer(_PopenFactory(page_count=2)).render(
            source, _sha256(source), symlink_output
        )
    with pytest.raises(PptxVisualRejected, match="page count"):
        _renderer(_PopenFactory(page_count=1)).render(
            source, _sha256(source), tmp_path / "wrong-pages"
        )


def test_render_enforces_image_and_total_byte_limits(tmp_path: Path) -> None:
    source = tmp_path / "lesson.pptx"
    _pptx(source)

    with pytest.raises(PptxVisualRejected, match="image byte limit"):
        _renderer(
            _PopenFactory(page_count=2),
            max_image_bytes=32,
            max_total_bytes=64,
        ).render(source, _sha256(source), tmp_path / "visuals")


def test_render_timeout_terminates_the_process_group(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "lesson.pptx"
    _pptx(source)
    signals: list[tuple[int, int]] = []
    monkeypatch.setattr(
        "open_notebook.course.pptx_visual_renderer.os.killpg",
        lambda pid, signal_number: signals.append((pid, signal_number)),
    )

    with pytest.raises(PptxVisualUnavailable, match="timed out"):
        _renderer(_PopenFactory(page_count=2, timeout=True), timeout_seconds=0.01).render(
            source, _sha256(source), tmp_path / "visuals"
        )

    assert signals


def test_missing_converter_is_explicitly_unavailable(tmp_path: Path) -> None:
    source = tmp_path / "lesson.pptx"
    _pptx(source)

    with pytest.raises(PptxVisualUnavailable, match="LibreOffice"):
        PptxVisualRenderer(converter=None, converter_locator=lambda: None).render(
            source, _sha256(source), tmp_path / "visuals"
        )


@pytest.mark.skipif(
    os.environ.get("OPEN_NOTEBOOK_RUN_REAL_PPTX_VISUAL_SMOKE") != "1",
    reason="set OPEN_NOTEBOOK_RUN_REAL_PPTX_VISUAL_SMOKE=1 for real LibreOffice smoke",
)
def test_real_libreoffice_pdfium_smoke(tmp_path: Path) -> None:
    source = GOLD_ROOT / "stem-evidence-gold.pptx"
    renderer = PptxVisualRenderer()
    if renderer.locate_converter() is None:
        pytest.skip("LibreOffice soffice is genuinely unavailable")

    rendered = renderer.render(source, _sha256(source), tmp_path / "visuals")

    assert list(rendered) == [1, 2, 3]
    assert all(path.read_bytes().startswith(b"\x89PNG") for path in rendered.values())
    assert len({path.read_bytes() for path in rendered.values()}) == 3
