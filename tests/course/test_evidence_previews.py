"""Safe, Course-owned evidence preview and original navigation contracts."""

from __future__ import annotations

import hashlib
import io
import json
from pathlib import Path
from types import SimpleNamespace
from typing import Literal
from unittest.mock import AsyncMock

import pytest
from fastapi.testclient import TestClient

from api.course_service import CourseService
from open_notebook.course.evidence_service import EvidenceInputError, EvidenceService
from open_notebook.course.models import Course
from open_notebook.course.pptx_visual_renderer import (
    PptxVisualRejected,
    PptxVisualUnavailable,
)
from open_notebook.domain.notebook import Asset, Source
from open_notebook.exceptions import NotFoundError


def _pptx(path: Path, slide_count: int = 3) -> None:
    from pptx import Presentation

    presentation = Presentation()
    for index in range(1, slide_count + 1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(10, 10, 500, 100)
        box.text = f"Slide {index} grounded text"
    presentation.save(str(path))


def _pdf(path: Path) -> None:
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    with path.open("wb") as handle:
        writer.write(handle)


def _png_bytes(width: int = 1280, height: int = 720) -> bytes:
    from PIL import Image

    stream = io.BytesIO()
    Image.new("RGB", (width, height), "white").save(stream, format="PNG")
    return stream.getvalue()


def _course(source_ids: list[str] | None = None) -> Course:
    selected = source_ids if source_ids is not None else ["source:one"]
    return Course(
        id="course:one",
        title="Course",
        notebook="notebook:one",
        source_ids=selected,
        primary_source_ids=selected,
    )


def _preview_records() -> list[
    tuple[int, str, str, tuple[float, float, float, float] | None]
]:
    return [
        (1, "#/texts/0", "<script>alert('x')</script> & grounded\x00", None),
        (1, "#/texts/1", "Second provenance-backed line", None),
        (2, "#/texts/2", "A" * 20_000, None),
    ]


def _anchor(
    service: EvidenceService,
    *,
    source_hash: str,
    preview_path: str | None,
    visual_preview_path: str | None = None,
    visual_preview_status: Literal["available", "text_only"] = "text_only",
    index: int = 1,
):
    return service.make_anchor(
        course_id="course:one",
        source_id="source:one",
        source_sha256=source_hash,
        kind="pptx_slide",
        index=index,
        block_key="#/texts/0",
        quote="Grounded slide text",
        source_role="PRIMARY",
        preview_path=preview_path,
        visual_preview_path=visual_preview_path,
        visual_preview_status=visual_preview_status,
    )


def test_pptx_previews_are_deterministic_escaped_bounded_and_per_slide(
    tmp_path: Path,
) -> None:
    service = EvidenceService(data_root=tmp_path / "cache", allowed_roots=[tmp_path])
    source_hash = hashlib.sha256(b"slides").hexdigest()

    first = service.write_pptx_previews(
        course_id="course:one",
        source_id="source:one",
        source_sha256=source_hash,
        records=_preview_records(),
        slide_count=3,
    )
    second = service.write_pptx_previews(
        course_id="course:one",
        source_id="source:one",
        source_sha256=source_hash,
        records=_preview_records(),
        slide_count=3,
    )

    assert first == second
    assert set(first) == {1, 2, 3}
    for relative_path in first.values():
        assert not Path(relative_path).is_absolute()
        content = (service.data_root / relative_path).read_bytes()
        assert len(content) <= service.MAX_PREVIEW_BYTES
        assert content.count(b"<tspan") <= service.MAX_PREVIEW_LINES
        assert b"\x00" not in content

    first_svg = (service.data_root / first[1]).read_text(encoding="utf-8")
    assert "&lt;script&gt;" in first_svg
    assert "&amp; grounded" in first_svg
    assert "<script" not in first_svg.lower()
    assert "javascript:" not in first_svg.lower()
    assert "href=" not in first_svg.lower()
    assert "onload=" not in first_svg.lower()


def test_preview_loader_verifies_exact_identity_and_content_hash(tmp_path: Path) -> None:
    service = EvidenceService(data_root=tmp_path / "cache", allowed_roots=[tmp_path])
    source_hash = hashlib.sha256(b"slides").hexdigest()
    previews = service.write_pptx_previews(
        course_id="course:one",
        source_id="source:one",
        source_sha256=source_hash,
        records=_preview_records(),
        slide_count=2,
    )
    anchor = _anchor(
        service, source_hash=source_hash, preview_path=previews[1]
    )

    asset = service.load_preview_asset(
        anchor, course_id="course:one", source_hash=source_hash
    )

    assert asset.filename == "slide-0001.svg"
    assert asset.content.startswith(b"<svg")

    path = service.data_root / previews[1]
    path.write_bytes(path.read_bytes() + b"tampered")
    with pytest.raises(EvidenceInputError, match="identity|hash"):
        service.load_preview_asset(
            anchor, course_id="course:one", source_hash=source_hash
        )


def test_preview_loader_prefers_hash_bound_visual_png_and_rejects_tampering(
    tmp_path: Path,
) -> None:
    service = EvidenceService(data_root=tmp_path / "cache", allowed_roots=[tmp_path])
    source_hash = hashlib.sha256(b"slides").hexdigest()
    previews = service.write_pptx_previews(
        course_id="course:one",
        source_id="source:one",
        source_sha256=source_hash,
        records=_preview_records(),
        slide_count=2,
    )
    png = _png_bytes()
    digest = hashlib.sha256(png).hexdigest()[:16]
    relative_dir = service._preview_directory(
        "course:one", "source:one", source_hash
    )
    relative_visual = relative_dir / f"slide-0001-{digest}.png"
    visual_path = service.data_root / relative_visual
    visual_path.write_bytes(png)
    anchor = _anchor(
        service,
        source_hash=source_hash,
        preview_path=previews[1],
        visual_preview_path=relative_visual.as_posix(),
        visual_preview_status="available",
    )

    asset = service.load_preview_asset(
        anchor, course_id="course:one", source_hash=source_hash
    )

    assert asset.content == png
    assert asset.filename == "slide-0001.png"
    assert asset.media_type == "image/png"
    assert asset.mode == "visual"

    visual_path.write_bytes(png + b"tampered")
    with pytest.raises(EvidenceInputError, match="identity|hash"):
        service.load_preview_asset(
            anchor, course_id="course:one", source_hash=source_hash
        )


def test_visual_preview_rejects_invalid_dimensions_and_symlink(tmp_path: Path) -> None:
    service = EvidenceService(data_root=tmp_path / "cache", allowed_roots=[tmp_path])
    source_hash = hashlib.sha256(b"slides").hexdigest()
    previews = service.write_pptx_previews(
        course_id="course:one",
        source_id="source:one",
        source_sha256=source_hash,
        records=_preview_records(),
        slide_count=2,
    )
    invalid = _png_bytes(width=640, height=360)
    digest = hashlib.sha256(invalid).hexdigest()[:16]
    relative_dir = service._preview_directory(
        "course:one", "source:one", source_hash
    )
    relative_visual = relative_dir / f"slide-0001-{digest}.png"
    visual_path = service.data_root / relative_visual
    visual_path.write_bytes(invalid)
    anchor = _anchor(
        service,
        source_hash=source_hash,
        preview_path=previews[1],
        visual_preview_path=relative_visual.as_posix(),
        visual_preview_status="available",
    )

    with pytest.raises(EvidenceInputError, match="dimension"):
        service.load_preview_asset(
            anchor, course_id="course:one", source_hash=source_hash
        )

    outside = tmp_path / "outside.png"
    outside.write_bytes(_png_bytes())
    visual_path.unlink()
    visual_path.symlink_to(outside)
    with pytest.raises(EvidenceInputError, match="symbolic"):
        service.load_preview_asset(
            anchor, course_id="course:one", source_hash=source_hash
        )


@pytest.mark.parametrize(
    ("preview_path", "message"),
    [
        ("../../outside.svg", "identity|path"),
        (None, "unavailable|missing"),
    ],
)
def test_preview_loader_rejects_traversal_and_missing_identity(
    tmp_path: Path, preview_path: str | None, message: str
) -> None:
    service = EvidenceService(data_root=tmp_path / "cache", allowed_roots=[tmp_path])
    source_hash = hashlib.sha256(b"slides").hexdigest()
    anchor = _anchor(
        service, source_hash=source_hash, preview_path=preview_path
    )

    with pytest.raises(EvidenceInputError, match=message):
        service.load_preview_asset(
            anchor, course_id="course:one", source_hash=source_hash
        )


def test_preview_loader_rejects_missing_file_and_symlink(tmp_path: Path) -> None:
    service = EvidenceService(data_root=tmp_path / "cache", allowed_roots=[tmp_path])
    source_hash = hashlib.sha256(b"slides").hexdigest()
    previews = service.write_pptx_previews(
        course_id="course:one",
        source_id="source:one",
        source_sha256=source_hash,
        records=_preview_records(),
        slide_count=2,
    )
    anchor = _anchor(
        service, source_hash=source_hash, preview_path=previews[1]
    )
    preview_file = service.data_root / previews[1]
    content = preview_file.read_bytes()
    preview_file.unlink()

    with pytest.raises(EvidenceInputError, match="missing"):
        service.load_preview_asset(
            anchor, course_id="course:one", source_hash=source_hash
        )

    outside = tmp_path / "outside.svg"
    outside.write_bytes(content)
    preview_file.symlink_to(outside)
    with pytest.raises(EvidenceInputError, match="symbolic"):
        service.load_preview_asset(
            anchor, course_id="course:one", source_hash=source_hash
        )


@pytest.mark.asyncio
async def test_pptx_build_assigns_one_internal_preview_identity_per_slide(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source_path = tmp_path / "slides.pptx"
    _pptx(source_path)
    service = EvidenceService(data_root=tmp_path / "cache", allowed_roots=[tmp_path])
    payload = {
        "pages": {
            str(index): {"size": {"width": 1280, "height": 720}}
            for index in range(1, 4)
        },
        "texts": [
            {
                "self_ref": f"#/texts/{index - 1}",
                "text": f"Slide {index} grounded text",
                "prov": [{"page_no": index}],
            }
            for index in range(1, 4)
        ],
    }
    monkeypatch.setattr(Course, "get", AsyncMock(return_value=_course()))
    monkeypatch.setattr(
        Source,
        "get",
        AsyncMock(
            return_value=Source(
                id="source:one",
                title="Slides",
                asset=Asset(file_path=str(source_path)),
            )
        ),
    )
    monkeypatch.setattr(
        service, "_extract_docling_content", AsyncMock(return_value=json.dumps(payload))
    )
    monkeypatch.setattr(
        service,
        "_persist",
        AsyncMock(side_effect=lambda **kwargs: kwargs["anchors"]),
    )

    anchors = await service.build(
        course_id="course:one",
        source_id="source:one",
        source_role="PRIMARY",
    )

    assert [anchor.locator.index for anchor in anchors] == [1, 2, 3]
    assert all(anchor.preview_path for anchor in anchors)
    assert len({anchor.preview_path for anchor in anchors}) == 3
    assert all(
        (service.data_root / str(anchor.preview_path)).is_file()
        for anchor in anchors
    )


class _VisualRenderer:
    def __init__(self, outcome: str = "available") -> None:
        self.outcome = outcome
        self.calls: list[tuple[Path, str, Path]] = []

    def render(
        self, path: Path, expected_sha256: str, output_dir: Path
    ) -> dict[int, Path]:
        self.calls.append((path, expected_sha256, output_dir))
        if self.outcome == "unavailable":
            raise PptxVisualUnavailable("soffice missing")
        if self.outcome == "rejected":
            raise PptxVisualRejected("unsafe embedded object")
        output_dir.mkdir(parents=True, exist_ok=True)
        rendered: dict[int, Path] = {}
        for index in range(1, 4):
            target = output_dir / f"slide-{index:04d}.png"
            target.write_bytes(_png_bytes())
            rendered[index] = target
        return rendered


def _docling_payload() -> str:
    return json.dumps({
        "pages": {
            str(index): {"size": {"width": 1280, "height": 720}}
            for index in range(1, 4)
        },
        "texts": [
            {
                "self_ref": f"#/texts/{index - 1}",
                "text": f"Slide {index} grounded text",
                "prov": [{"page_no": index, "bbox": {"l": 100, "t": 100, "r": 500, "b": 300}}],
            }
            for index in range(1, 4)
        ],
    })


async def _build_with_renderer(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    renderer: _VisualRenderer,
):
    source_path = tmp_path / "slides.pptx"
    _pptx(source_path)
    service = EvidenceService(
        data_root=tmp_path / "cache",
        allowed_roots=[tmp_path],
        visual_renderer=renderer,
    )
    monkeypatch.setattr(Course, "get", AsyncMock(return_value=_course()))
    monkeypatch.setattr(
        Source,
        "get",
        AsyncMock(return_value=Source(
            id="source:one",
            title="Slides",
            asset=Asset(file_path=str(source_path)),
        )),
    )
    monkeypatch.setattr(
        service, "_extract_docling_content", AsyncMock(return_value=_docling_payload())
    )
    monkeypatch.setattr(
        service,
        "_persist",
        AsyncMock(side_effect=lambda **kwargs: kwargs["anchors"]),
    )
    return service, await service.build(
        course_id="course:one",
        source_id="source:one",
        source_role="PRIMARY",
    )


@pytest.mark.asyncio
async def test_pptx_build_persists_visual_and_text_preview_tracks(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    renderer = _VisualRenderer()

    service, anchors = await _build_with_renderer(tmp_path, monkeypatch, renderer)

    assert len(renderer.calls) == 1
    assert all(anchor.preview_path for anchor in anchors)
    assert all(anchor.visual_preview_status == "available" for anchor in anchors)
    assert all(anchor.visual_preview_path for anchor in anchors)
    assert len({anchor.visual_preview_path for anchor in anchors}) == 3
    assert all(
        not Path(str(anchor.visual_preview_path)).is_absolute()
        and (service.data_root / str(anchor.visual_preview_path)).is_file()
        for anchor in anchors
    )


@pytest.mark.asyncio
async def test_pptx_build_records_honest_text_only_fallback(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    service, anchors = await _build_with_renderer(
        tmp_path, monkeypatch, _VisualRenderer("unavailable")
    )

    assert all(anchor.preview_path for anchor in anchors)
    assert all(anchor.visual_preview_status == "text_only" for anchor in anchors)
    assert all(anchor.visual_preview_path is None for anchor in anchors)
    assert all(
        (service.data_root / str(anchor.preview_path)).is_file()
        for anchor in anchors
    )


@pytest.mark.asyncio
async def test_pptx_build_fails_closed_for_unsafe_visual_input(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    with pytest.raises(EvidenceInputError, match="unsafe"):
        await _build_with_renderer(
            tmp_path, monkeypatch, _VisualRenderer("rejected")
        )


@pytest.mark.asyncio
async def test_course_preview_service_checks_course_ownership_and_current_hash(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source_path = tmp_path / "slides.pptx"
    _pptx(source_path)
    service = EvidenceService(data_root=tmp_path / "cache", allowed_roots=[tmp_path])
    source_hash = service.sha256_file(source_path)
    previews = service.write_pptx_previews(
        course_id="course:one",
        source_id="source:one",
        source_sha256=source_hash,
        records=_preview_records(),
        slide_count=3,
    )
    anchor = _anchor(
        service, source_hash=source_hash, preview_path=previews[1]
    )
    anchor.id = "course_evidence_anchor:one"
    source = Source(
        id="source:one",
        title="Slides",
        asset=Asset(file_path=str(source_path)),
    )
    monkeypatch.setattr(CourseService, "get_course", AsyncMock(return_value=_course()))
    query = AsyncMock(return_value=[anchor.model_dump(mode="json")])
    monkeypatch.setattr("api.course_service.repo_query", query)
    monkeypatch.setattr(Source, "get", AsyncMock(return_value=source))
    monkeypatch.setattr("api.course_service.EvidenceService", lambda: service)

    loaded = await CourseService.get_evidence_preview("course:one", anchor.anchor_id)

    assert loaded.content.startswith(b"<svg")
    assert query.await_args is not None
    assert query.await_args.args[1]["anchor_ids"] == [anchor.anchor_id]

    monkeypatch.setattr(
        CourseService, "get_course", AsyncMock(return_value=_course([]))
    )
    with pytest.raises(NotFoundError, match="not found"):
        await CourseService.get_evidence_preview("course:one", anchor.anchor_id)

    monkeypatch.setattr(CourseService, "get_course", AsyncMock(return_value=_course()))
    _pptx(source_path, slide_count=4)
    with pytest.raises(EvidenceInputError, match="changed"):
        await CourseService.get_evidence_preview("course:one", anchor.anchor_id)


@pytest.mark.asyncio
async def test_course_source_service_returns_owned_unchanged_pdf(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source_path = tmp_path / "lesson.pdf"
    _pdf(source_path)
    service = EvidenceService(data_root=tmp_path / "cache", allowed_roots=[tmp_path])
    source_hash = service.sha256_file(source_path)
    anchor = service.make_anchor(
        course_id="course:one",
        source_id="source:one",
        source_sha256=source_hash,
        kind="pdf_page",
        index=1,
        block_key="#/texts/0",
        quote="Grounded page text",
        source_role="PRIMARY",
    )
    anchor.id = "course_evidence_anchor:pdf"
    monkeypatch.setattr(CourseService, "get_course", AsyncMock(return_value=_course()))
    monkeypatch.setattr(
        "api.course_service.repo_query",
        AsyncMock(return_value=[anchor.model_dump(mode="json")]),
    )
    monkeypatch.setattr(
        Source,
        "get",
        AsyncMock(
            return_value=Source(
                id="source:one",
                title="Lesson",
                asset=Asset(file_path=str(source_path)),
            )
        ),
    )
    monkeypatch.setattr("api.course_service.EvidenceService", lambda: service)

    asset = await CourseService.get_evidence_source("course:one", anchor.anchor_id)

    assert asset.path == source_path.resolve()
    assert asset.filename == "lesson.pdf"
    assert asset.kind == "pdf"


def test_course_preview_endpoint_sets_visual_mode_media_and_security_headers(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setattr(
        CourseService,
        "get_evidence_preview",
        AsyncMock(
            side_effect=[
                SimpleNamespace(
                    content=_png_bytes(),
                    filename="slide-0001.png",
                    media_type="image/png",
                    mode="visual",
                ),
                SimpleNamespace(
                    content=b'<svg xmlns="http://www.w3.org/2000/svg"></svg>',
                    filename="slide-0001.svg",
                    media_type="image/svg+xml",
                    mode="text_only",
                ),
            ]
        ),
        raising=False,
    )
    from api.main import app

    client = TestClient(app)
    visual = client.get(
        "/api/courses/course:one/evidence/anchors/anchor:one/preview"
    )
    fallback = client.get(
        "/api/courses/course:one/evidence/anchors/anchor:two/preview"
    )

    assert visual.status_code == 200
    assert visual.headers["content-type"].startswith("image/png")
    assert visual.headers["x-course-preview-mode"] == "visual"
    assert visual.headers["x-content-type-options"] == "nosniff"
    assert visual.headers["content-disposition"].startswith("inline")
    assert visual.content.startswith(b"\x89PNG")
    assert fallback.status_code == 200
    assert fallback.headers["content-type"].startswith("image/svg+xml")
    assert fallback.headers["x-course-preview-mode"] == "text_only"
    assert "default-src 'none'" in fallback.headers["content-security-policy"]
    assert fallback.content.startswith(b"<svg")


def test_course_source_endpoint_uses_inline_pdf_and_attachment_pptx(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    pdf = tmp_path / "lesson.pdf"
    pptx = tmp_path / "slides.pptx"
    pdf.write_bytes(b"%PDF-test")
    pptx.write_bytes(b"pptx-test")
    source = AsyncMock(
        side_effect=[
            SimpleNamespace(path=pdf, filename="lesson.pdf", kind="pdf"),
            SimpleNamespace(path=pptx, filename="slides.pptx", kind="pptx"),
        ]
    )
    monkeypatch.setattr(
        CourseService, "get_evidence_source", source, raising=False
    )
    from api.main import app

    client = TestClient(app)
    inline = client.get(
        "/api/courses/course:one/evidence/anchors/anchor:pdf/source"
    )
    download = client.get(
        "/api/courses/course:one/evidence/anchors/anchor:pptx/source?download=true"
    )

    assert inline.status_code == 200
    assert inline.headers["content-type"].startswith("application/pdf")
    assert inline.headers["content-disposition"].startswith("inline")
    assert download.status_code == 200
    assert download.headers["content-disposition"].startswith("attachment")
    assert source.await_count == 2
