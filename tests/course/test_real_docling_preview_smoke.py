"""Opt-in real Docling smoke for synthetic PDF/PPTX evidence previews."""

from __future__ import annotations

import os
from pathlib import Path

import pytest

from open_notebook.course.evidence_service import EvidenceService

pytestmark = pytest.mark.skipif(
    os.getenv("OPEN_NOTEBOOK_RUN_REAL_DOCLING_SMOKE") != "1",
    reason="set OPEN_NOTEBOOK_RUN_REAL_DOCLING_SMOKE=1 for the local runtime gate",
)


def _synthetic_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    for index in range(1, 4):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        box.text = f"Synthetic vector lesson slide {index}: v{index} = {index} m/s"
    presentation.save(str(path))


def _synthetic_scanned_pdf(path: Path) -> None:
    from PIL import Image, ImageDraw

    pages = []
    for index in range(1, 3):
        image = Image.new("RGB", (1200, 800), "white")
        drawing = ImageDraw.Draw(image)
        drawing.text(
            (80, 100),
            f"Synthetic projectile page {index}\nposition equals velocity times time",
            fill="black",
            spacing=24,
        )
        pages.append(image)
    pages[0].save(path, "PDF", save_all=True, append_images=pages[1:], resolution=144)


def test_real_docling_extracts_three_slides_and_two_pdf_pages_with_previews(
    tmp_path: Path,
) -> None:
    pptx_path = tmp_path / "synthetic-vectors.pptx"
    pdf_path = tmp_path / "synthetic-projectile.pdf"
    _synthetic_pptx(pptx_path)
    _synthetic_scanned_pdf(pdf_path)
    service = EvidenceService(data_root=tmp_path / "course-evidence", allowed_roots=[tmp_path])

    pptx_records = service.docling_records(
        service._extract_docling_sync(pptx_path, "pptx"), "pptx"
    )
    assert {record[0] for record in pptx_records} == {1, 2, 3}
    previews = service.write_pptx_previews(
        course_id="course:smoke",
        source_id="source:pptx",
        source_sha256=service.sha256_file(pptx_path),
        records=pptx_records,
        slide_count=3,
    )
    assert set(previews) == {1, 2, 3}
    assert all((service.data_root / relative).is_file() for relative in previews.values())

    pdf_records = service.docling_records(
        service._extract_docling_sync(pdf_path, "pdf"), "pdf"
    )
    assert {record[0] for record in pdf_records} == {1, 2}
