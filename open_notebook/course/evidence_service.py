"""Grounded, immutable Course evidence extraction with Docling provenance."""

from __future__ import annotations

import asyncio
import hashlib
import inspect
import io
import json
import os
import re
import tempfile
import textwrap
import zipfile
from dataclasses import dataclass
from html import escape
from pathlib import Path
from typing import Any, Literal, Protocol, cast
from urllib.parse import urlparse

from open_notebook.config import UPLOADS_FOLDER
from open_notebook.database.repository import ensure_record_id, repo_query
from open_notebook.domain.notebook import Source
from open_notebook.exceptions import ConfigurationError, InvalidInputError

from .locking import course_job_lock
from .models import Course, CourseEvidenceAnchor
from .pptx_visual_renderer import (
    PptxVisualRejected,
    PptxVisualRenderer,
    PptxVisualUnavailable,
)
from .v2_contracts import EvidenceCategory, EvidenceClassification

EvidenceKind = Literal["pdf", "pptx"]
SourceRole = Literal["PRIMARY", "SUPPLEMENT"]
DoclingRecord = tuple[int, str, str, tuple[float, float, float, float] | None]


class VisualRenderer(Protocol):
    def render(
        self, path: Path, expected_sha256: str, output_dir: Path
    ) -> dict[int, Path]: ...


@dataclass(frozen=True)
class EvidencePreviewAsset:
    content: bytes
    filename: str
    media_type: Literal["image/png", "image/svg+xml"]
    mode: Literal["visual", "text_only"]


@dataclass(frozen=True)
class EvidenceSourceAsset:
    path: Path
    filename: str
    kind: EvidenceKind


class EvidenceInputError(InvalidInputError, ValueError):
    """A permanent, actionable source-file or evidence-integrity failure."""


class EvidenceConfigurationError(ConfigurationError):
    """The local evidence runtime is not installed or configured."""


class EvidenceService:
    """Build and persist anchors from a Course-owned Source asset."""

    MAX_SOURCE_BYTES = 100 * 1024 * 1024
    MAX_PPTX_MEMBERS = 10_000
    MAX_PPTX_EXPANDED_BYTES = 500 * 1024 * 1024
    MAX_PREVIEW_LINES = 12
    MAX_PREVIEW_LINE_CHARS = 96
    MAX_PREVIEW_TEXT_CHARS = MAX_PREVIEW_LINES * MAX_PREVIEW_LINE_CHARS
    MAX_PREVIEW_BYTES = 64 * 1024
    MAX_VISUAL_PREVIEW_BYTES = PptxVisualRenderer.DEFAULT_MAX_IMAGE_BYTES
    VISUAL_PREVIEW_WIDTH = PptxVisualRenderer.OUTPUT_WIDTH
    MAX_VISUAL_PREVIEW_HEIGHT = PptxVisualRenderer.MAX_OUTPUT_HEIGHT
    _PNG_SIGNATURE = b"\x89PNG\r\n\x1a\n"
    _ASSESSMENT_LABELS: tuple[tuple[EvidenceCategory, re.Pattern[str]], ...] = (
        (
            "prerequisite",
            re.compile(r"\b(?:prerequisite|recall|review)\b|先修|复习", re.I),
        ),
        (
            "worked_example",
            re.compile(r"\b(?:worked[ -]?example|example)\b|例题|示例", re.I),
        ),
        ("answer", re.compile(r"\b(?:answer|solution)\b|答案|解答", re.I)),
        (
            "exercise",
            re.compile(
                r"\b(?:exercise|problem|practice|question)\b|练习|习题|问题", re.I
            ),
        ),
        (
            "theorem",
            re.compile(
                r"\b(?:theorem|lemma|proposition|corollary)\b|定理|引理|命题|推论", re.I
            ),
        ),
        ("definition", re.compile(r"\bdefinition\b|定义", re.I)),
        ("figure", re.compile(r"\b(?:figure|table|diagram|chart)\b|图|表", re.I)),
    )
    _SOURCE_NUMBER_SUFFIX = re.compile(
        r"[\s:_#]*(?:no\.?|编号)?\s*([0-9]+(?:\.[0-9]+){0,5}[a-z]?)",
        re.I,
    )
    _ANSWER_SOURCE_NUMBER_SUFFIX = re.compile(
        r"[\s:_#]*(?:to\s+exercise\s+)?(?:no\.?|编号)?\s*"
        r"([0-9]+(?:\.[0-9]+){0,5}[a-z]?)",
        re.I,
    )

    @classmethod
    def _number_after_label(
        cls, category: EvidenceCategory, text: str, label: re.Match[str]
    ) -> re.Match[str] | None:
        pattern = (
            cls._ANSWER_SOURCE_NUMBER_SUFFIX
            if category == "answer"
            else cls._SOURCE_NUMBER_SUFFIX
        )
        return pattern.match(text[label.end() :])

    def __init__(
        self,
        data_root: Path | None = None,
        allowed_roots: list[Path] | None = None,
        model_root: Path | None = None,
        visual_renderer: VisualRenderer | None = None,
    ) -> None:
        self.data_root = (data_root or Path("notebook_data/course_evidence")).resolve()
        self.model_root = (
            model_root or self.data_root.parent / "course_models"
        ).resolve()
        roots = allowed_roots if allowed_roots is not None else [Path(UPLOADS_FOLDER)]
        self.allowed_roots = [root.resolve() for root in roots]
        self.visual_renderer = visual_renderer or PptxVisualRenderer()

    @staticmethod
    def normalize_text(value: str) -> str:
        return re.sub(r"\s+", " ", value or "").strip()

    @classmethod
    def classify_assessment_anchor(
        cls, anchor: CourseEvidenceAnchor
    ) -> EvidenceClassification:
        """Classify an anchor without trusting model-supplied labels."""

        quote = cls.normalize_text(anchor.locator.quote)
        block_key = cls.normalize_text(anchor.locator.block_key)
        category: EvidenceCategory = "unclassified"
        confidence = "low"
        block_matches = [
            (match.start(), candidate, match)
            for candidate, pattern in cls._ASSESSMENT_LABELS
            if (match := pattern.search(block_key)) is not None
        ]
        selected_match: re.Match[str] | None = None
        selected_text = ""
        if block_matches:
            _, category, selected_match = min(block_matches, key=lambda item: item[0])
            selected_text = block_key
            confidence = "high"
        quote_matches = [
            (match.start(), candidate, match)
            for candidate, pattern in cls._ASSESSMENT_LABELS
            if (match := pattern.search(quote)) is not None
        ]
        if not block_matches:
            if quote_matches:
                first_position, category, selected_match = min(
                    quote_matches, key=lambda item: item[0]
                )
                selected_text = quote
                confidence = "high" if first_position == 0 else "medium"
        number_match = (
            cls._number_after_label(category, selected_text, selected_match)
            if selected_match is not None
            else None
        )
        if number_match is None and block_matches:
            leading_quote_match = next(
                (
                    match
                    for position, candidate, match in quote_matches
                    if position == 0 and candidate == category
                ),
                None,
            )
            if leading_quote_match is not None:
                number_match = cls._number_after_label(
                    category, quote, leading_quote_match
                )
        return EvidenceClassification(
            anchor_id=anchor.anchor_id,
            category=category,
            confidence=confidence,
            source_number=number_match.group(1) if number_match else None,
        )

    @classmethod
    def assessment_context(
        cls,
        anchor: CourseEvidenceAnchor,
        classification: EvidenceClassification,
    ) -> str:
        if classification.anchor_id != anchor.anchor_id:
            raise EvidenceInputError(
                "Evidence classification does not match its anchor."
            )
        source_number = classification.source_number or "none"
        quote = json.dumps(cls.normalize_text(anchor.locator.quote), ensure_ascii=False)
        return (
            f"[anchor_id={anchor.anchor_id} category={classification.category} "
            f"confidence={classification.confidence} source_number={source_number} "
            f"role={anchor.source_role}] {quote}"
        )

    @classmethod
    def quote_sha256(cls, quote: str) -> str:
        return hashlib.sha256(cls.normalize_text(quote).encode("utf-8")).hexdigest()

    _quote_hash = quote_sha256

    @staticmethod
    def deterministic_anchor_id(
        *,
        course_id: str,
        source_id: str,
        source_sha256: str,
        kind: str,
        index: int,
        block_key: str,
        quote_sha256: str,
    ) -> str:
        stable_key = "|".join(
            [
                course_id,
                source_id,
                source_sha256,
                kind,
                str(index),
                block_key,
                quote_sha256,
            ]
        )
        return f"anchor:{hashlib.sha256(stable_key.encode()).hexdigest()[:32]}"

    @staticmethod
    def sha256_file(file_path: Path) -> str:
        digest = hashlib.sha256()
        with file_path.open("rb") as handle:
            for chunk in iter(lambda: handle.read(1024 * 1024), b""):
                digest.update(chunk)
        return digest.hexdigest()

    @staticmethod
    def validate_extension(file_path: str | Path) -> EvidenceKind:
        raw = str(file_path)
        if urlparse(raw).scheme.lower() in {"http", "https"}:
            raise EvidenceInputError(
                "URL sources are not supported; attach a local PDF or PPTX file."
            )
        suffix = Path(raw).suffix.lower()
        if suffix == ".pdf":
            return "pdf"
        if suffix == ".pptx":
            return "pptx"
        if suffix == ".ppt":
            raise EvidenceInputError(
                "Legacy .ppt is not supported; convert the file to PPTX first."
            )
        raise EvidenceInputError("Course sources must be PDF or PPTX files.")

    def resolve_safe_source_path(self, file_path: str | Path) -> Path:
        candidate = Path(file_path)
        if candidate.is_symlink():
            raise EvidenceInputError("Course source must not be a symbolic link.")
        try:
            resolved = candidate.resolve(strict=True)
        except (FileNotFoundError, OSError) as exc:
            raise EvidenceInputError("Course source file is missing.") from exc
        if not resolved.is_file():
            raise EvidenceInputError("Course source path must be a regular file.")
        if not any(
            resolved == root or root in resolved.parents for root in self.allowed_roots
        ):
            raise EvidenceInputError(
                "Course source must be inside an allowed source directory."
            )
        size = resolved.stat().st_size
        if size == 0:
            raise EvidenceInputError("Course source file is empty.")
        if size > self.MAX_SOURCE_BYTES:
            raise EvidenceInputError("Course source exceeds the 100 MB size limit.")
        return resolved

    @staticmethod
    def _validate_file(path: Path, kind: EvidenceKind) -> None:
        if kind == "pdf":
            try:
                from pypdf import PdfReader

                reader = PdfReader(str(path), strict=True)
                if reader.is_encrypted:
                    raise EvidenceInputError(
                        "Encrypted PDF files are not supported; remove the password first."
                    )
                if len(reader.pages) == 0:
                    raise EvidenceInputError("The PDF contains no readable pages.")
                for page in reader.pages:
                    _ = page.mediabox
            except EvidenceInputError:
                raise
            except Exception as exc:
                raise EvidenceInputError(
                    "The PDF is corrupt or cannot be read; export it again."
                ) from exc
            return
        try:
            with zipfile.ZipFile(path) as archive:
                members = archive.infolist()
                if (
                    len(members) > EvidenceService.MAX_PPTX_MEMBERS
                    or sum(member.file_size for member in members)
                    > EvidenceService.MAX_PPTX_EXPANDED_BYTES
                ):
                    raise EvidenceInputError(
                        "The PPTX expands beyond the safe validation limit."
                    )
                names = {member.filename for member in members}
                if "[Content_Types].xml" not in names or not any(
                    name.startswith("ppt/slides/slide") for name in names
                ):
                    raise EvidenceInputError(
                        "The PPTX is corrupt or contains no readable slides."
                    )
                bad_member = archive.testzip()
                if bad_member is not None:
                    raise EvidenceInputError(
                        f"The PPTX is corrupt near {Path(bad_member).name}."
                    )
            # Opening the package validates its OPC relationships and XML,
            # unlike a filename-only ZIP check. This is intentionally local
            # structural validation: it does not invoke Docling or OCR.
            from pptx import Presentation

            presentation = Presentation(str(path))
            if len(presentation.slides) == 0:
                raise EvidenceInputError("The PPTX contains no readable slides.")
            for slide in presentation.slides:
                _ = len(slide.shapes)
        except EvidenceInputError:
            raise
        except Exception as exc:
            raise EvidenceInputError(
                "The PPTX is corrupt or cannot be read; export it again."
            ) from exc

    def validate_local_source_file(
        self, file_path: str | Path
    ) -> tuple[Path, EvidenceKind]:
        """Validate a local original without running extraction or OCR."""

        path = self.resolve_safe_source_path(file_path)
        kind = self.validate_extension(path)
        self._validate_file(path, kind)
        return path, kind

    @staticmethod
    def _serialize_docling_document(document: Any) -> str:
        export_json = getattr(document, "export_to_json", None)
        if callable(export_json):
            value = export_json()
            return value if isinstance(value, str) else json.dumps(value)
        export_dict = getattr(document, "export_to_dict", None)
        if callable(export_dict):
            return json.dumps(export_dict(), ensure_ascii=False)
        raise EvidenceInputError(
            "Docling returned an unsupported document; update the Docling runtime."
        )

    def _extract_docling_sync(self, path: Path, kind: EvidenceKind) -> str:
        try:
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.pipeline_options import (
                PdfPipelineOptions,
                RapidOcrOptions,
            )
            from docling.document_converter import DocumentConverter, PdfFormatOption
        except ImportError as exc:
            raise EvidenceConfigurationError(
                "Docling runtime is unavailable. Install the project with the Docling extra."
            ) from exc

        if kind == "pdf":
            options = PdfPipelineOptions(
                do_ocr=True,
                ocr_options=RapidOcrOptions(backend="onnxruntime"),
                do_formula_enrichment=True,
                do_picture_description=False,
                do_chart_extraction=False,
                artifacts_path=self.model_root if self.model_root.is_dir() else None,
            )
            converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=options)
                }
            )
        else:
            converter = DocumentConverter()
        try:
            result = converter.convert(str(path))
            return self._serialize_docling_document(result.document)
        except (EvidenceInputError, EvidenceConfigurationError):
            raise
        except Exception as exc:
            raise EvidenceInputError(
                "Docling evidence extraction failed; verify the file and local model cache."
            ) from exc

    async def _extract_docling_content(self, path: Path, kind: EvidenceKind) -> str:
        return await asyncio.to_thread(self._extract_docling_sync, path, kind)

    @classmethod
    def docling_records(
        cls, raw_content: str, kind: EvidenceKind
    ) -> list[DoclingRecord]:
        del kind
        try:
            payload = json.loads(raw_content)
        except (TypeError, json.JSONDecodeError):
            return []
        if not isinstance(payload, dict):
            return []
        pages = payload.get("pages", {})
        texts = payload.get("texts", [])
        if not isinstance(texts, list):
            return []
        records: list[DoclingRecord] = []
        for item_index, item in enumerate(texts, start=1):
            if not isinstance(item, dict):
                continue
            quote = item.get("text") or item.get("orig") or ""
            provenance = item.get("prov") or []
            if (
                not isinstance(quote, str)
                or not quote.strip()
                or not isinstance(provenance, list)
            ):
                continue
            for prov_index, prov in enumerate(provenance, start=1):
                if not isinstance(prov, dict):
                    continue
                page_no = prov.get("page_no")
                if not isinstance(page_no, int) or page_no < 1:
                    continue
                page = pages.get(str(page_no), {}) if isinstance(pages, dict) else {}
                size = page.get("size", {}) if isinstance(page, dict) else {}
                width = size.get("width") if isinstance(size, dict) else None
                height = size.get("height") if isinstance(size, dict) else None
                raw_bbox = prov.get("bbox")
                bbox: tuple[float, float, float, float] | None = None
                if isinstance(raw_bbox, dict):
                    values = [raw_bbox.get(key) for key in ("l", "t", "r", "b")]
                    origin = str(raw_bbox.get("coord_origin", "TOPLEFT")).upper()
                elif isinstance(raw_bbox, (list, tuple)) and len(raw_bbox) == 4:
                    values = list(raw_bbox)
                    origin = "TOPLEFT"
                else:
                    values = []
                    origin = "TOPLEFT"
                if (
                    len(values) == 4
                    and all(isinstance(value, (int, float)) for value in values)
                    and isinstance(width, (int, float))
                    and isinstance(height, (int, float))
                    and width > 0
                    and height > 0
                ):
                    numeric = cast(
                        tuple[float | int, float | int, float | int, float | int],
                        tuple(values),
                    )
                    left, top, right, bottom = (float(value) for value in numeric)
                    if origin == "BOTTOMLEFT":
                        top, bottom = height - top, height - bottom
                    bbox = tuple(
                        max(0.0, min(1.0, point))
                        for point in (
                            left / width,
                            top / height,
                            right / width,
                            bottom / height,
                        )
                    )  # type: ignore[assignment]
                records.append(
                    (
                        page_no,
                        str(
                            item.get("self_ref")
                            or f"docling-text-{item_index}-prov-{prov_index}"
                        ),
                        cls.normalize_text(quote)[:4000],
                        bbox,
                    )
                )
        return records

    _docling_records = docling_records

    def make_anchor(
        self,
        *,
        course_id: str,
        source_id: str,
        source_sha256: str,
        kind: Literal["pdf_page", "pptx_slide"],
        index: int,
        block_key: str,
        quote: str,
        source_role: SourceRole,
        bbox: tuple[float, float, float, float] | None = None,
        preview_path: str | None = None,
        visual_preview_path: str | None = None,
        visual_preview_status: Literal["available", "text_only"] = "text_only",
    ) -> CourseEvidenceAnchor:
        normalized_quote = self.normalize_text(quote)
        quote_hash = self.quote_sha256(normalized_quote)
        anchor_id = self.deterministic_anchor_id(
            course_id=course_id,
            source_id=source_id,
            source_sha256=source_sha256,
            kind=kind,
            index=index,
            block_key=block_key,
            quote_sha256=quote_hash,
        )
        return CourseEvidenceAnchor(
            course=course_id,
            source=source_id,
            anchor_id=anchor_id,
            locator={
                "source_id": source_id,
                "kind": kind,
                "index": index,
                "block_key": block_key,
                "quote": normalized_quote,
                "content_sha256": source_sha256,
                "bbox": bbox,
            },
            quote_sha256=quote_hash,
            source_role=source_role,
            preview_path=preview_path,
            visual_preview_path=visual_preview_path,
            visual_preview_status=visual_preview_status,
        )

    def manifest_path(self, course_id: str, source_id: str, source_sha256: str) -> Path:
        course_namespace = hashlib.sha256(course_id.encode("utf-8")).hexdigest()
        safe_name = hashlib.sha256(source_id.encode("utf-8")).hexdigest() + ".json"
        return self.data_root / course_namespace / source_sha256 / safe_name

    @staticmethod
    def _course_namespace(course_id: str) -> str:
        return hashlib.sha256(course_id.encode("utf-8")).hexdigest()

    @staticmethod
    def _source_namespace(source_id: str) -> str:
        return hashlib.sha256(source_id.encode("utf-8")).hexdigest()

    def _preview_directory(
        self, course_id: str, source_id: str, source_sha256: str
    ) -> Path:
        return (
            Path(self._course_namespace(course_id))
            / source_sha256
            / "previews"
            / self._source_namespace(source_id)
        )

    def _ensure_cache_directory(self, relative: Path) -> Path:
        if relative.is_absolute() or ".." in relative.parts:
            raise EvidenceInputError("Evidence preview cache path is invalid.")
        self.data_root.mkdir(parents=True, exist_ok=True)
        current = self.data_root
        for part in relative.parts:
            current = current / part
            if current.exists():
                if current.is_symlink() or not current.is_dir():
                    raise EvidenceInputError(
                        "Evidence preview cache must not contain symbolic links."
                    )
            else:
                current.mkdir()
        resolved = current.resolve(strict=True)
        if self.data_root != resolved and self.data_root not in resolved.parents:
            raise EvidenceInputError("Evidence preview cache path escaped its root.")
        return current

    @classmethod
    def render_slide_preview(cls, slide_index: int, quotes: list[str]) -> bytes:
        """Render bounded provenance text into a static, non-executable SVG."""

        xml_safe = "".join(
            character
            for character in " ".join(quotes)
            if ord(character) in {0x09, 0x0A, 0x0D}
            or 0x20 <= ord(character) <= 0xD7FF
            or 0xE000 <= ord(character) <= 0xFFFD
            or 0x10000 <= ord(character) <= 0x10FFFF
        )
        normalized = cls.normalize_text(xml_safe)[: cls.MAX_PREVIEW_TEXT_CHARS]
        lines = textwrap.wrap(
            normalized,
            width=cls.MAX_PREVIEW_LINE_CHARS,
            break_long_words=True,
            break_on_hyphens=False,
        )[: cls.MAX_PREVIEW_LINES]
        if not lines:
            lines = ["No provenance-backed text was extracted for this slide."]
        tspans = "".join(
            f'<tspan x="72" dy="{42 if position else 0}">{escape(line)}</tspan>'
            for position, line in enumerate(lines)
        )
        svg = (
            '<svg xmlns="http://www.w3.org/2000/svg" width="1280" height="720" '
            'viewBox="0 0 1280 720" role="img" aria-labelledby="title desc">'
            '<title id="title">Course evidence slide preview</title>'
            f'<desc id="desc">Slide {slide_index} provenance text</desc>'
            '<rect width="1280" height="720" fill="#f8fafc"/>'
            '<rect x="32" y="32" width="1216" height="656" rx="18" '
            'fill="#ffffff" stroke="#cbd5e1" stroke-width="2"/>'
            f'<text x="72" y="92" font-family="system-ui, sans-serif" '
            f'font-size="30" font-weight="700" fill="#0f172a">Slide {slide_index}</text>'
            '<text x="72" y="146" font-family="system-ui, sans-serif" '
            f'font-size="24" fill="#334155">{tspans}</text>'
            "</svg>"
        ).encode("utf-8")
        if len(svg) > cls.MAX_PREVIEW_BYTES:
            raise EvidenceInputError("Evidence preview exceeds its safe size limit.")
        return svg

    def write_pptx_previews(
        self,
        *,
        course_id: str,
        source_id: str,
        source_sha256: str,
        records: list[DoclingRecord],
        slide_count: int,
    ) -> dict[int, str]:
        """Create one deterministic cache-relative SVG identity per slide."""

        if slide_count < 1 or slide_count > self.MAX_PPTX_MEMBERS:
            raise EvidenceInputError("The PPTX slide count is outside safe limits.")
        quotes_by_slide: dict[int, list[str]] = {
            index: [] for index in range(1, slide_count + 1)
        }
        for index, _block_key, quote, _bbox in records:
            if index in quotes_by_slide:
                quotes_by_slide[index].append(quote)
        relative_directory = self._preview_directory(
            course_id, source_id, source_sha256
        )
        directory = self._ensure_cache_directory(relative_directory)
        previews: dict[int, str] = {}
        for index in range(1, slide_count + 1):
            content = self.render_slide_preview(index, quotes_by_slide[index])
            digest = hashlib.sha256(content).hexdigest()[:16]
            filename = f"slide-{index:04d}-{digest}.svg"
            target = directory / filename
            if target.is_symlink():
                raise EvidenceInputError(
                    "Evidence preview cache must not contain symbolic links."
                )
            target.write_bytes(content)
            previews[index] = (relative_directory / filename).as_posix()
        return previews

    @classmethod
    def _validate_png_content(cls, content: bytes) -> tuple[int, int]:
        if (
            len(content) < 24
            or len(content) > cls.MAX_VISUAL_PREVIEW_BYTES
            or not content.startswith(cls._PNG_SIGNATURE)
        ):
            raise EvidenceInputError("Evidence visual preview is not a bounded PNG.")
        width = int.from_bytes(content[16:20], "big")
        height = int.from_bytes(content[20:24], "big")
        if (
            width != cls.VISUAL_PREVIEW_WIDTH
            or height <= 0
            or height > cls.MAX_VISUAL_PREVIEW_HEIGHT
        ):
            raise EvidenceInputError("Evidence visual preview dimensions are invalid.")
        try:
            from PIL import Image

            with Image.open(io.BytesIO(content)) as image:
                if image.format != "PNG" or image.size != (width, height):
                    raise EvidenceInputError(
                        "Evidence visual preview dimensions are invalid."
                    )
                image.verify()
        except EvidenceInputError:
            raise
        except Exception as exc:
            raise EvidenceInputError(
                "Evidence visual preview is not a valid PNG."
            ) from exc
        return width, height

    def write_pptx_visual_previews(
        self,
        *,
        path: Path,
        course_id: str,
        source_id: str,
        source_sha256: str,
        slide_count: int,
    ) -> dict[int, str]:
        """Render and cache one content-addressed, cache-relative PNG per slide."""

        if slide_count < 1 or slide_count > self.MAX_PPTX_MEMBERS:
            raise EvidenceInputError("The PPTX slide count is outside safe limits.")
        relative_directory = self._preview_directory(
            course_id, source_id, source_sha256
        )
        cache_directory = self._ensure_cache_directory(relative_directory)
        with tempfile.TemporaryDirectory(prefix="course-evidence-visual-") as raw_temp:
            stage_directory = Path(raw_temp)
            rendered = self.visual_renderer.render(
                path, source_sha256, stage_directory
            )
            expected_indices = set(range(1, slide_count + 1))
            if set(rendered) != expected_indices:
                raise PptxVisualRejected(
                    "Rendered slide identities do not match the PPTX."
                )
            if self.sha256_file(path) != source_sha256:
                raise PptxVisualRejected(
                    "PPTX source changed during visual rendering."
                )

            stage_root = stage_directory.resolve(strict=True)
            previews: dict[int, str] = {}
            for index in range(1, slide_count + 1):
                rendered_path = Path(rendered[index])
                if rendered_path.is_symlink():
                    raise PptxVisualRejected(
                        "Rendered visual preview must not be a symbolic link."
                    )
                try:
                    resolved_rendered = rendered_path.resolve(strict=True)
                except (FileNotFoundError, OSError) as exc:
                    raise PptxVisualRejected(
                        "Rendered visual preview is missing."
                    ) from exc
                if (
                    stage_root not in resolved_rendered.parents
                    or not resolved_rendered.is_file()
                ):
                    raise PptxVisualRejected(
                        "Rendered visual preview escaped its staging directory."
                    )
                if (
                    resolved_rendered.stat().st_size <= 0
                    or resolved_rendered.stat().st_size
                    > self.MAX_VISUAL_PREVIEW_BYTES
                ):
                    raise PptxVisualRejected(
                        "Rendered visual preview exceeds its byte limit."
                    )
                content = resolved_rendered.read_bytes()
                try:
                    self._validate_png_content(content)
                except EvidenceInputError as exc:
                    raise PptxVisualRejected(str(exc)) from exc
                digest = hashlib.sha256(content).hexdigest()[:16]
                filename = f"slide-{index:04d}-{digest}.png"
                target = cache_directory / filename
                if target.is_symlink():
                    raise PptxVisualRejected(
                        "Evidence visual cache must not contain symbolic links."
                    )
                if target.exists():
                    if (
                        not target.is_file()
                        or target.stat().st_size != len(content)
                        or target.stat().st_size > self.MAX_VISUAL_PREVIEW_BYTES
                        or target.read_bytes() != content
                    ):
                        raise PptxVisualRejected(
                            "Evidence visual preview identity collision."
                        )
                else:
                    temporary = cache_directory / f".{filename}.tmp-{os.getpid()}"
                    try:
                        if temporary.exists() or temporary.is_symlink():
                            raise PptxVisualRejected(
                                "Evidence visual cache staging path is unsafe."
                            )
                        temporary.write_bytes(content)
                        os.replace(temporary, target)
                    finally:
                        if temporary.exists() and not temporary.is_symlink():
                            temporary.unlink()
                previews[index] = (relative_directory / filename).as_posix()
            return previews

    def _read_preview_content(
        self,
        *,
        stored_path: str,
        expected_directory: Path,
        filename_pattern: str,
        max_bytes: int,
        label: str,
    ) -> bytes:
        relative = Path(stored_path)
        filename_match = re.fullmatch(filename_pattern, relative.name)
        if (
            relative.is_absolute()
            or ".." in relative.parts
            or relative.as_posix() != stored_path
            or relative.parent != expected_directory
            or filename_match is None
        ):
            raise EvidenceInputError(f"Evidence {label} identity or path is invalid.")
        candidate = self.data_root
        for part in relative.parts:
            candidate = candidate / part
            if candidate.is_symlink():
                raise EvidenceInputError(
                    f"Evidence {label} must not be a symbolic link."
                )
        try:
            resolved = candidate.resolve(strict=True)
        except (FileNotFoundError, OSError) as exc:
            raise EvidenceInputError(f"Evidence {label} file is missing.") from exc
        if self.data_root != resolved and self.data_root not in resolved.parents:
            raise EvidenceInputError(
                f"Evidence {label} path escaped its cache root."
            )
        if not resolved.is_file():
            raise EvidenceInputError(f"Evidence {label} file is missing.")
        if resolved.stat().st_size <= 0 or resolved.stat().st_size > max_bytes:
            raise EvidenceInputError(f"Evidence {label} file has an invalid size.")
        content = resolved.read_bytes()
        if hashlib.sha256(content).hexdigest()[:16] != filename_match.group(1):
            raise EvidenceInputError(
                f"Evidence {label} identity hash does not match."
            )
        return content

    def load_preview_asset(
        self,
        anchor: CourseEvidenceAnchor,
        *,
        course_id: str,
        source_hash: str,
    ) -> EvidencePreviewAsset:
        """Load only the exact immutable preview identity stored on an anchor."""

        self.validate_anchor_integrity(
            anchor, course_id=course_id, source_hash=source_hash
        )
        if anchor.locator.kind != "pptx_slide" or not anchor.preview_path:
            raise EvidenceInputError("Evidence preview is unavailable for this anchor.")
        expected_directory = self._preview_directory(
            course_id, anchor.source, source_hash
        )
        if anchor.visual_preview_status == "available":
            if not anchor.visual_preview_path:
                raise EvidenceInputError(
                    "Evidence visual preview identity is missing."
                )
            content = self._read_preview_content(
                stored_path=anchor.visual_preview_path,
                expected_directory=expected_directory,
                filename_pattern=(
                    rf"slide-{anchor.locator.index:04d}-([0-9a-f]{{16}})\.png"
                ),
                max_bytes=self.MAX_VISUAL_PREVIEW_BYTES,
                label="visual preview",
            )
            self._validate_png_content(content)
            return EvidencePreviewAsset(
                content=content,
                filename=f"slide-{anchor.locator.index:04d}.png",
                media_type="image/png",
                mode="visual",
            )

        content = self._read_preview_content(
            stored_path=anchor.preview_path,
            expected_directory=expected_directory,
            filename_pattern=(
                rf"slide-{anchor.locator.index:04d}-([0-9a-f]{{16}})\.svg"
            ),
            max_bytes=self.MAX_PREVIEW_BYTES,
            label="preview",
        )
        lowered = content.lower()
        if not content.startswith(b'<svg xmlns="http://www.w3.org/2000/svg"') or any(
            token in lowered
            for token in (
                b"<script",
                b"javascript:",
                b"<foreignobject",
                b" href=",
                b" xlink:href=",
                b" onload=",
                b" onclick=",
            )
        ):
            raise EvidenceInputError("Evidence preview content is not safe SVG.")
        return EvidencePreviewAsset(
            content=content,
            filename=f"slide-{anchor.locator.index:04d}.svg",
            media_type="image/svg+xml",
            mode="text_only",
        )

    @staticmethod
    def pptx_slide_count(path: Path) -> int:
        from pptx import Presentation

        return len(Presentation(str(path)).slides)

    @staticmethod
    def manifest(
        *,
        course_id: str,
        source_id: str,
        source_sha256: str,
        anchors: list[CourseEvidenceAnchor],
    ) -> dict[str, Any]:
        return {
            "version": 1,
            "course_id": course_id,
            "source_id": source_id,
            "source_sha256": source_sha256,
            "anchors": [anchor.model_dump(mode="json") for anchor in anchors],
        }

    @classmethod
    def validate_anchor_integrity(
        cls,
        anchor: CourseEvidenceAnchor,
        *,
        course_id: str,
        source_hash: str,
    ) -> None:
        if anchor.course != course_id:
            raise EvidenceInputError(
                f"Evidence anchor {anchor.anchor_id} does not belong to this Course."
            )
        if not anchor.is_current:
            raise EvidenceInputError(f"Evidence anchor {anchor.anchor_id} is stale.")
        if anchor.locator.source_id != anchor.source:
            raise EvidenceInputError(
                f"Evidence anchor {anchor.anchor_id} has a source mismatch."
            )
        if anchor.locator.content_sha256 != source_hash:
            raise EvidenceInputError(
                f"Evidence source for {anchor.anchor_id} changed; rebuild evidence."
            )
        if cls.quote_sha256(anchor.locator.quote) != anchor.quote_sha256:
            raise EvidenceInputError(
                f"Evidence anchor {anchor.anchor_id} has a quote hash mismatch."
            )
        expected_anchor_id = cls.deterministic_anchor_id(
            course_id=course_id,
            source_id=anchor.source,
            source_sha256=source_hash,
            kind=anchor.locator.kind,
            index=anchor.locator.index,
            block_key=anchor.locator.block_key,
            quote_sha256=anchor.quote_sha256,
        )
        if anchor.anchor_id != expected_anchor_id:
            raise EvidenceInputError(
                f"Evidence anchor {anchor.anchor_id} has an anchor ID mismatch."
            )

    @classmethod
    def retrieval_context(
        cls,
        anchors: list[CourseEvidenceAnchor],
        *,
        selected_anchor_ids: list[str],
        course_id: str,
        source_hashes: dict[str, str],
    ) -> list[str]:
        by_id = {anchor.anchor_id: anchor for anchor in anchors}
        context: list[str] = []
        for anchor_id in selected_anchor_ids:
            anchor = by_id.get(anchor_id)
            if anchor is None:
                raise EvidenceInputError(f"Unknown evidence anchor: {anchor_id}")
            source_hash = source_hashes.get(anchor.source)
            if source_hash is None:
                raise EvidenceInputError(
                    f"Evidence source {anchor.source} is not selected for this Course."
                )
            cls.validate_anchor_integrity(
                anchor, course_id=course_id, source_hash=source_hash
            )
            locator = anchor.locator
            context.append(
                f"{anchor.source_role} {locator.kind} {locator.index} "
                f"[{anchor.anchor_id}]: {locator.quote}"
            )
        return context

    @staticmethod
    def _assert_role(course: Course, source_id: str, role: SourceRole) -> None:
        in_primary = source_id in course.primary_source_ids
        in_supplement = source_id in course.supplement_source_ids
        actual = (
            "PRIMARY"
            if in_primary and not in_supplement
            else "SUPPLEMENT"
            if in_supplement and not in_primary
            else None
        )
        if actual != role:
            detail = actual or "no role"
            raise EvidenceInputError(
                f"Source is associated with this Course as {detail}, not {role}."
            )

    async def _persist(
        self,
        *,
        course_id: str,
        source_id: str,
        source_role: SourceRole,
        source_hash: str,
        kind: EvidenceKind,
        anchors: list[CourseEvidenceAnchor],
    ) -> list[CourseEvidenceAnchor]:
        statement = """
BEGIN TRANSACTION;
LET $evidence = (
    UPSERT ONLY evidence
    SET course = $course_id,
        source = $source_id,
        title = $source_title,
        kind = $kind,
        file_hash = $source_hash,
        source_hash = $source_hash,
        source_role = $source_role,
        status = 'ready'
    WHERE course = $course_id AND source = $source_id
    RETURN AFTER
);
UPDATE course_evidence_anchor
SET is_current = false
WHERE course = $course_id
  AND source = $source_id
  AND anchor_id NOT IN $anchor_ids;
FOR $anchor IN $anchors {
    UPSERT course_evidence_anchor
    SET course = $course_id,
        source = $source_id,
        evidence = $evidence.id,
        anchor_id = $anchor.anchor_id,
        locator = $anchor.locator,
        quote_sha256 = $anchor.quote_sha256,
        source_role = $anchor.source_role,
        preview_path = $anchor.preview_path,
        visual_preview_path = $anchor.visual_preview_path,
        visual_preview_status = $anchor.visual_preview_status,
        is_current = true
    WHERE course = $course_id
      AND source = $source_id
      AND anchor_id = $anchor.anchor_id;
};
COMMIT TRANSACTION;
"""
        await repo_query(
            statement,
            {
                "course_id": ensure_record_id(course_id),
                "source_id": ensure_record_id(source_id),
                "source_title": source_id,
                "kind": kind,
                "source_hash": source_hash,
                "source_role": source_role,
                "anchor_ids": [anchor.anchor_id for anchor in anchors],
                "anchors": [
                    {
                        "anchor_id": anchor.anchor_id,
                        "locator": anchor.locator.model_dump(mode="json"),
                        "quote_sha256": anchor.quote_sha256,
                        "source_role": anchor.source_role,
                        "preview_path": anchor.preview_path,
                        "visual_preview_path": anchor.visual_preview_path,
                        "visual_preview_status": anchor.visual_preview_status,
                    }
                    for anchor in anchors
                ],
            },
        )
        return anchors

    async def build(
        self,
        *,
        course_id: str,
        source_id: str,
        source_role: SourceRole,
    ) -> list[CourseEvidenceAnchor]:
        course = await Course.get(course_id)
        self._assert_role(course, source_id, source_role)
        source = await Source.get(source_id)
        file_path = source.asset.file_path if source.asset else None
        if not file_path:
            raise EvidenceInputError(
                "Course Source has no local asset file; upload a PDF or PPTX."
            )
        path, kind = self.validate_local_source_file(file_path)
        source_hash = self.sha256_file(path)

        async with course_job_lock():
            extracted = self._extract_docling_content(path, kind)
            raw_content = (
                await extracted if inspect.isawaitable(extracted) else extracted
            )
            records = self.docling_records(str(raw_content), kind)
            if not records:
                raise EvidenceInputError(
                    "Docling returned no provenance-backed text; evidence was not created."
                )
            locator_kind: Literal["pdf_page", "pptx_slide"] = (
                "pdf_page" if kind == "pdf" else "pptx_slide"
            )
            slide_count = self.pptx_slide_count(path) if kind == "pptx" else 0
            preview_paths = (
                self.write_pptx_previews(
                    course_id=course_id,
                    source_id=source_id,
                    source_sha256=source_hash,
                    records=records,
                    slide_count=slide_count,
                )
                if kind == "pptx"
                else {}
            )
            visual_preview_paths: dict[int, str] = {}
            visual_preview_status: Literal["available", "text_only"] = "text_only"
            if kind == "pptx":
                try:
                    visual_preview_paths = await asyncio.to_thread(
                        self.write_pptx_visual_previews,
                        path=path,
                        course_id=course_id,
                        source_id=source_id,
                        source_sha256=source_hash,
                        slide_count=slide_count,
                    )
                except PptxVisualUnavailable:
                    visual_preview_paths = {}
                except PptxVisualRejected as exc:
                    raise EvidenceInputError(str(exc)) from exc
                else:
                    visual_preview_status = "available"
            if self.sha256_file(path) != source_hash:
                raise EvidenceInputError(
                    "Course source changed during evidence extraction; rebuild it."
                )
            anchors = [
                self.make_anchor(
                    course_id=course_id,
                    source_id=source_id,
                    source_sha256=source_hash,
                    kind=locator_kind,
                    index=index,
                    block_key=block_key,
                    quote=quote,
                    source_role=source_role,
                    bbox=bbox,
                    preview_path=preview_paths.get(index),
                    visual_preview_path=visual_preview_paths.get(index),
                    visual_preview_status=visual_preview_status,
                )
                for index, block_key, quote, bbox in records
            ]
            manifest_path = self.manifest_path(course_id, source_id, source_hash)
            manifest_path.parent.mkdir(parents=True, exist_ok=True)
            manifest_path.write_text(
                json.dumps(
                    self.manifest(
                        course_id=course_id,
                        source_id=source_id,
                        source_sha256=source_hash,
                        anchors=anchors,
                    ),
                    ensure_ascii=False,
                    indent=2,
                    sort_keys=True,
                ),
                encoding="utf-8",
            )
            return await self._persist(
                course_id=course_id,
                source_id=source_id,
                source_role=source_role,
                source_hash=source_hash,
                kind=kind,
                anchors=anchors,
            )
